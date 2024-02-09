#Importing Modules

import speech_recognition as sr
import pyttsx3
import pywhatkit
import datetime
import wikipedia
import pyjokes
import random 
import requests
import os
import time
from pptx import Presentation

listener = sr.Recognizer()
engine = pyttsx3.init()
voices = engine.getProperty('voices')
engine.setProperty('voice', voices[1].id)

#assistant talking to us
def talk(text): 
    engine.say(text)
    engine.runAndWait()

#giving commands to assistant
def take_command():
    try:
        with sr.Microphone() as source:
            print('I am listening...')
            voice = listener.listen(source)
            command = listener.recognize_google(voice)
            command = command.lower()
            if 'seema aunty' in command:
                command = command.replace('seema aunty', '')
                print(command)
    except:
        pass
    return command

NEWS_API_KEY = "a2da9d62c2834040bb22643a24e3ca07"

#getting latest news
def get_news():
    try:
        url = f"https://newsapi.org/v2/top-headlines?country=us&apiKey={NEWS_API_KEY}"
        response = requests.get(url)
        news_data = response.json()

        if news_data.get("status") == "ok":
            articles = news_data.get("articles")
            return articles

    except Exception as e:
        print("Error fetching news:", str(e))
    
    return []

# Get your API key from Merriam-Webster
MERRIAM_WEBSTER_API_KEY = "489e40b0-2861-407c-afe6-10377eece25d" 

#getting definitions
def get_word_definition(word):
    try:
        url = f"https://www.dictionaryapi.com/api/v3/references/collegiate/json/{word}?key={MERRIAM_WEBSTER_API_KEY}"
        response = requests.get(url)
        definitions = response.json()

        if definitions:
            # Extract the first definition from the response
            definition = definitions[0]
            if isinstance(definition, str):
                return definition
            elif "shortdef" in definition:
                return ", ".join(definition["shortdef"])
        else:
            return "Sorry, I couldn't find the definition for that word."

    except Exception as e:
        print("Error fetching word definition:", str(e))
        return "Sorry, an error occurred while fetching the word definition."

#Information about places 
def get_place_info(place_name):
    try:
        page = wikipedia.page(place_name)
        summary = wikipedia.summary(place_name)
        return summary

    except wikipedia.exceptions.DisambiguationError as e:
        suggestions = e.options
        return f"Did you mean {', '.join(suggestions)}?"
    
    except wikipedia.exceptions.PageError:
        return "Sorry, I couldn't find information about that place."

    except Exception as e:
        print("Error fetching place information:", str(e))
        return "Sorry, an error occurred while fetching information about the place."

#Opening the presentation
def open_powerpoint_presentation(presentation_path):
    try:
        os.startfile(presentation_path)
    except Exception as e:
        print("Error opening PowerPoint presentation:", str(e))
        talk("Sorry, I couldn't open the PowerPoint presentation.")

#Reading the contents of the presenatation 
def read_slide_content(slide):
    content = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            content += shape.text + "\n"
    return content

def read_presentation(presentation_path):
    try:
        presentation = Presentation(presentation_path)
        for slide in presentation.slides:
            content = read_slide_content(slide)
            talk(content)
            time.sleep(2)  # Adjust the duration for each slide

    except Exception as e:
        print("Error reading PowerPoint presentation:", str(e))
        talk("Sorry, I couldn't read the presentation.")


# List of interesting facts
facts = [
    "Did you know that RCB has never won an IPL trophy till now. Ha-Ha-Ha!"
    "Did you know that honey never spoils?",
    "Did you know that the Eiffel Tower can be 15 cm taller during the summer?",
    "Did you know that a group of flamingos is called a 'flamboyance'?",
    "Did you know that a day on Venus is longer than its year?",
    "Did you know that there are more possible iterations of a game of chess than there are atoms in the known universe?",
    "Did you know that Apples float on water!",
    "Did you know that the world wastes about 1 billion metric tons of food each year.",
    "Did you know that the shortest war in the history lasted just for 38 to 45 minutes between Britain and Zanzibar in 1896",
    "Did you know that honeybees can recognize human faces",
    "Did you know that he world's largest desert is not the Sahara, but Antarctica. Deserts are defined by their low precipitation, not just sand.",
    "Did you know that Astronauts cannot burp in space due to the absence of gravity.",
    "Did you know that The world's oldest known living tree is a Great Basin bristlecone pine named Methuselah, estimated to be over 4,800 years old.",
    "Did you know that Octopuses have three hearts and blue blood.",
    "Did you know that the Great Wall of China is not visible from space without aid, contrary to popular belief.",
    "Did you know that the first recorded game of baseball was played in 1846 in Hoboken, New Jersey.",
    "Did you know that the world's smallest mammal is the bumblebee bat, with a wingspan of about 1.1 inches.",]

#generating random facts
def get_random_fact():
    return random.choice(facts)

def run_alexa():
    command = take_command()
    print(command)
   
    if 'play' in command:
        song = command.replace('play', '')
        talk('playing ' + song)
        pywhatkit.playonyt(song)

    elif 'time' in command:
        time = datetime.datetime.now().strftime('%I:%M %p')
        talk('Current time is ' + time)
   
    elif 'who is' in command:
        person = command.replace('who the heck is', '')
        info = wikipedia.summary(person, 1)
        print(info)
        talk(info)
  
    elif 'where is' in command:
        country= command.replace('where the heck is', '')
        info = wikipedia.summary(country, 1)
        print(info)
        talk(info)
   
    elif 'date' in command:
        talk('sorry, I have a headache')
    
    elif 'are you single' in command:
        talk('I am in a relationship with my dear wifi')
   
    elif 'joke' in command:
        talk(pyjokes.get_joke())
   
    elif 'girlfriend' in command :
        talk('Please shut Up! First go and see your face in the mirror and then ask me that question, ')
   
    elif 'fuck' in command:
        talk("Please dont allow me to use my powers, which may badly affect you emotionally, Teri maa ki chooth!")
   
    elif "fact" in command:
        facts= get_random_fact()
        talk( facts)

    elif "good morning" in command:
        talk("Good morning homo sapiens")
   
    elif "team" in command:
        talk("The Team members for the mini-project using python includes 2 students from CSE branch 1st year. Let's introduce Chethan P, USN: 1DT22CS040. Hemanth Kumar M, USN:1DT22CS057")
    
    elif "python professor" in command :
        talk("Introduction to python Programming for CSE-1 is taken by Vijayalakshmi Inamdar. You guys are very lucky to learn this beautiful language from a very beautiful teacher. She always tries to give her best to teach the concepts and modules. Thank you ma'am, i, SEEMA AUNTY, is very lucky to present in front of you. Hope you give my inventors full marks, in all upcoming tests and practicals exams.")
   
    elif "who are you" in command :
        talk('''Hello there! I am your digital companion, your helpful guide through the virtual realm. Allow me to introduce myself - I am SEEMA AUNTY, your very own virtual assistant, here to make your life easier, more entertaining, and a whole lot smarter.
    Just like a trusted friend. From answering your questions, providing the latest news, to telling you jokes and stories, there's hardly anything I can't do.
    I can understand your voice commands and respond in a way that feels like we're having a real conversation.
    But I'm not just about functionality – I also come with a touch of personality. You can ask me for a joke, a fun fact, or even share your thoughts, and I'll be here to engage and entertain you.
    So, whether you're looking for information, a bit of entertainment, or simply a friendly chat, don't hesitate to ask. Just say my name followed by your question or request, and I'll be right here, ready to serve.
    Let's embark on this exciting journey together and explore the endless possibilities I offer. Welcome to the world of SEEMA AUNTY!''')
    
    elif "thank you" in command :
        talk("Thank You for being here and showing your patience throughout the session. Okay, Byeeeee. see you guys again")

    elif "doubts" in command:
        talk("My Friends, Chethan and Hemanth, have given a very detailed and excellent presentation, about my working mechanism. If you still have any doubts,  then feel free to ask me anything! I will try to answer them. --Thank You!")
    
    elif "news" in command:
        articles = get_news()
        if articles:
            talk("Here are the top news headlines:")
            for index, article in enumerate(articles[:5], start=1):
                title = article.get("title")
                source = article.get("source").get("name")
                talk(f"Headline {index}: {title} from {source}")
        else:
            talk("Sorry, I couldn't fetch the latest news at the moment. Please try again later.")

    elif 'define' in command:
        word = command.replace('define', '').strip()
        definition = get_word_definition(word)
        talk(definition)

    elif 'tell me about' in command:
        place = command.replace('tell me about', '').strip()
        place_info = get_place_info(place)
        talk(place_info)

    elif 'tell me about' in command:
        person = command.replace('tell me about', '').strip()
        person_info = get_person_info(person)
        talk(person_info)

    elif 'open' in command:
        presentation_path = r'C:\Users\Hemanth Kumar M\Desktop\Virtual assistant.pptx' 
        open_powerpoint_presentation(presentation_path)

    elif 'read presentation' in command:
        presentation_path = r'C:\Users\Hemanth Kumar M\Desktop\Virtual assistant.pptx' 
        read_presentation(presentation_path)
    
    else:
        talk('Please say the command again.')


while True:
    run_alexa()
